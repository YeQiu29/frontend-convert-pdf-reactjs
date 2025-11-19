import os
import tempfile
import zipfile
from io import BytesIO
from typing import List

import uvicorn
from fastapi import FastAPI, File, UploadFile, HTTPException, Form
from fastapi.responses import FileResponse, StreamingResponse
from fastapi.middleware.cors import CORSMiddleware
from starlette.background import BackgroundTask

# --- Import library PDF ---
from pypdf import PdfWriter, PdfReader, PasswordType
from pdf2docx import Converter
from pdf2image import convert_from_bytes

# ... import lainnya ...
from PIL import Image

# --- Import untuk Watermark / Tanda Tangan ---
from reportlab.pdfgen import canvas
from reportlab.lib.units import inch
from reportlab.lib import colors
from reportlab.lib.pagesizes import A4
from reportlab.lib.utils import ImageReader # <-- [BARU] Tambahkan import ini

# --- Import untuk PowerPoint ---
from pptx import Presentation
from pptx.util import Inches

# --- Import untuk Excel ---
import camelot
import pandas as pd
import fitz  # [BARU] Ini adalah PyMuPDF
from openpyxl.utils.cell import get_column_letter # [BARU]
from openpyxl.drawing.image import Image as OpenPyXLImage # [BARU]

app = FastAPI(
    title="BigPDF Backend API",
    description="API untuk mengelola dan mengkonversi file PDF.",
    version="1.0.0"
)

# --- [TAMBAHKAN KODE INI] ---
# Ini adalah bagian yang penting untuk mengizinkan frontend Anda
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"], # Mengizinkan SEMUA origin (ganti dengan domain Anda saat produksi)
    allow_credentials=True,
    allow_methods=["*"], # Mengizinkan semua metode (POST, GET, dll)
    allow_headers=["*"], # Mengizinkan semua header
)
# --- [SELESAI TAMBAHAN] ---

# --- Helper Functions ---

def cleanup_file(path: str):
    try:
        os.remove(path)
    except Exception as e:
        print(f"Error cleaning up file {path}: {e}")

def cleanup_dir(path: str):
    import shutil
    try:
        shutil.rmtree(path)
    except Exception as e:
        print(f"Error cleaning up directory {path}: {e}")

def create_watermark_pdf(text: str) -> BytesIO:
    packet = BytesIO()
    can = canvas.Canvas(packet, pagesize=A4)
    width, height = A4
    can.setFont("Helvetica", 50)
    can.setFillAlpha(0.3)
    can.translate(width / 2, height / 2)
    can.rotate(45)
    can.drawCentredString(0, 0, text)
    can.save()
    packet.seek(0)
    return packet

def parse_page_range(range_str: str, max_pages: int) -> set:
    """Helper untuk mem-parsing string rentang halaman (cth: '1, 3, 5-7')"""
    pages = set()
    try:
        for part in range_str.split(','):
            part = part.strip()
            if '-' in part:
                start, end = part.split('-')
                start = int(start.strip())
                end = int(end.strip())
                if start < 1 or end > max_pages or start > end:
                    raise ValueError(f"Rentang '{part}' tidak valid.")
                # User (1-indexed) -> internal (0-indexed)
                for i in range(start, end + 1):
                    pages.add(i - 1)
            else:
                page = int(part.strip())
                if page < 1 or page > max_pages:
                    raise ValueError(f"Halaman '{page}' di luar rentang.")
                pages.add(page - 1) # 0-indexed
    except Exception as e:
        raise HTTPException(
            400, 
            f"Format rentang halaman tidak valid: '{range_str}'. Gunakan format seperti '1, 3, 5-7'. Error: {e}"
        )
    if not pages:
         raise HTTPException(400, "Tidak ada halaman yang dipilih.")
    return pages


# --- API Endpoints ---

@app.get("/")
def read_root():
    return {"message": "BigPDF API is running. Docs at /docs"}


@app.post("/merge", summary="Gabungkan beberapa PDF")
async def merge_pdfs(files: List[UploadFile] = File(..., description="File PDF yang akan digabung")):
    # (Kode tidak berubah)
    merger = PdfWriter()
    for file in files:
        if file.content_type != "application/pdf":
            raise HTTPException(400, "Hanya file PDF yang diizinkan.")
        pdf_bytes = await file.read()
        try:
            reader = PdfReader(BytesIO(pdf_bytes))
            if reader.is_encrypted:
                raise HTTPException(400, f"File {file.filename} terenkripsi. Harap buka sandi terlebih dahulu.")
            merger.append(BytesIO(pdf_bytes))
        except Exception as e:
            raise HTTPException(400, f"Error membaca {file.filename}: {e}")
    output_io = BytesIO()
    merger.write(output_io)
    merger.close()
    output_io.seek(0)
    return StreamingResponse(
        output_io,
        media_type="application/pdf",
        headers={"Content-Disposition": "attachment; filename=merged.pdf"}
    )


@app.post("/to-word", summary="Konversi PDF ke Word (.docx)")
async def pdf_to_word(file: UploadFile = File(..., description="File PDF yang akan dikonversi")):
    # (Kode tidak berubah)
    if file.content_type != "application/pdf":
        raise HTTPException(400, "Hanya file PDF yang diizinkan.")
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as temp_pdf:
            content = await file.read()
            temp_pdf.write(content)
            temp_pdf_path = temp_pdf.name
        output_path = tempfile.mktemp(suffix=".docx")
        cv = Converter(temp_pdf_path)
        cv.convert(output_path, start=0, end=None)
        cv.close()
        os.remove(temp_pdf_path)
        return FileResponse(
            path=output_path,
            media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            filename=f"{os.path.splitext(file.filename)[0]}.docx",
            background=BackgroundTask(cleanup_file, path=output_path)
        )
    except Exception as e:
        if 'temp_pdf_path' in locals() and os.path.exists(temp_pdf_path):
            os.remove(temp_pdf_path)
        if 'output_path' in locals() and os.path.exists(output_path):
            os.remove(output_path)
        raise HTTPException(500, f"Terjadi error saat konversi: {e}")


@app.post("/to-images", summary="Konversi PDF ke Gambar (ZIP)")
async def pdf_to_images(file: UploadFile = File(..., description="File PDF yang akan dikonversi")):
    # (Kode tidak berubah)
    if file.content_type != "application/pdf":
        raise HTTPException(400, "Hanya file PDF yang diizinkan.")
    try:
        temp_dir = tempfile.mkdtemp()
        pdf_bytes = await file.read()
        images = convert_from_bytes(pdf_bytes, output_folder=temp_dir, fmt='png')
        zip_io = BytesIO()
        with zipfile.ZipFile(zip_io, 'w') as zf:
            for i, image_path in enumerate(os.listdir(temp_dir)):
                full_path = os.path.join(temp_dir, image_path)
                new_filename = f"page_{i+1}.png"
                zf.write(full_path, arcname=new_filename)
        zip_io.seek(0)
        return StreamingResponse(
            zip_io,
            media_type="application/zip",
            headers={"Content-Disposition": f"attachment; filename={os.path.splitext(file.filename)[0]}.zip"},
            background=BackgroundTask(cleanup_dir, path=temp_dir)
        )
    except Exception as e:
        if 'temp_dir' in locals() and os.path.exists(temp_dir):
            cleanup_dir(temp_dir)
        raise HTTPException(500, f"Terjadi error saat konversi ke gambar: {e}. Pastikan Poppler terinstal.")


@app.post("/watermark", summary="Tambahkan watermark ke PDF")
async def add_watermark(
    file: UploadFile = File(..., description="File PDF utama."),
    text: str = Form(..., description="Teks untuk watermark.")
):
    # (Kode tidak berubah)
    if file.content_type != "application/pdf":
        raise HTTPException(400, "Hanya file PDF yang diizinkan.")
    try:
        pdf_bytes = await file.read()
        pdf_reader = PdfReader(BytesIO(pdf_bytes))
        if pdf_reader.is_encrypted:
            raise HTTPException(400, "File PDF terenkripsi. Buka sandi terlebih dahulu.")
        watermark_io = create_watermark_pdf(text)
        watermark_reader = PdfReader(watermark_io)
        watermark_page = watermark_reader.pages[0]
        writer = PdfWriter()
        for page in pdf_reader.pages:
            page.merge_page(watermark_page)
            writer.add_page(page)
        output_io = BytesIO()
        writer.write(output_io)
        output_io.seek(0)
        return StreamingResponse(
            output_io,
            media_type="application/pdf",
            headers={"Content-Disposition": f"attachment; filename=watermarked_{file.filename}"}
        )
    except Exception as e:
        raise HTTPException(500, f"Terjadi error: {e}")


@app.post("/lock", summary="Kunci PDF dengan sandi")
async def lock_pdf(
    file: UploadFile = File(..., description="File PDF yang akan dikunci."),
    password: str = Form(..., description="Sandi untuk PDF.")
):
    # (Kode tidak berubah)
    if file.content_type != "application/pdf":
        raise HTTPException(400, "Hanya file PDF yang diizinkan.")
    try:
        pdf_bytes = await file.read()
        reader = PdfReader(BytesIO(pdf_bytes))
        if reader.is_encrypted:
            raise HTTPException(400, "File sudah terenkripsi.")
        writer = PdfWriter()
        for page in reader.pages:
            writer.add_page(page)
        writer.encrypt(password)
        output_io = BytesIO()
        writer.write(output_io)
        output_io.seek(0)
        return StreamingResponse(
            output_io,
            media_type="application/pdf",
            headers={"Content-Disposition": f"attachment; filename=locked_{file.filename}"}
        )
    except Exception as e:
        raise HTTPException(500, f"Terjadi error: {e}")


@app.post("/unlock", summary="Hapus sandi dari PDF")
async def unlock_pdf(
    file: UploadFile = File(..., description="File PDF yang akan dibuka."),
    password: str = Form(..., description="Sandi PDF yang benar.")
):
    # (Kode perbaikan terakhir, tidak berubah)
    if file.content_type != "application/pdf":
        raise HTTPException(400, "Hanya file PDF yang diizinkan.")
    try:
        pdf_bytes = await file.read()
        reader = PdfReader(BytesIO(pdf_bytes))
        if not reader.is_encrypted:
            raise HTTPException(400, "File tidak terenkripsi.")
        
        result = reader.decrypt(password)
        if result == PasswordType.NOT_DECRYPTED:
             raise HTTPException(403, "Sandi salah.")

        writer = PdfWriter()
        for page in reader.pages:
            writer.add_page(page)
        output_io = BytesIO()
        writer.write(output_io)
        output_io.seek(0)
        return StreamingResponse(
            output_io,
            media_type="application/pdf",
            headers={"Content-Disposition": f"attachment; filename=unlocked_{file.filename}"}
        )
    except HTTPException as e:
        raise e
    except Exception as e:
        raise HTTPException(500, f"Gagal mendekripsi: {e}. Pastikan sandi benar.")


# --- ENDPOINT LAMA DIGANTI DENGAN YANG INI ---
@app.post("/split", summary="Pisahkan PDF berdasarkan rentang halaman")
async def split_pdf_flexible(
    file: UploadFile = File(..., description="File PDF yang akan dipisah."),
    page_range: str = Form(..., description="Halaman yang akan diekstrak (cth: '6' atau '1-3, 5')")
):
    """
    Ekstrak halaman berdasarkan rentang (cth: '1, 3, 5-7')
    dan kembalikan ZIP berisi 2 file: 'extracted.pdf' dan 'remaining.pdf'.
    """
    if file.content_type != "application/pdf":
        raise HTTPException(400, "Hanya file PDF yang diizinkan.")

    try:
        temp_dir = tempfile.mkdtemp()
        pdf_bytes = await file.read()
        reader = PdfReader(BytesIO(pdf_bytes))

        if reader.is_encrypted:
            raise HTTPException(400, "File PDF terenkripsi. Buka sandi terlebih dahulu.")
        
        total_pages = len(reader.pages)
        
        # Dapatkan set halaman (0-indexed) yang akan diekstrak
        extracted_indices = parse_page_range(page_range, total_pages)
        
        writer_extracted = PdfWriter()
        writer_remaining = PdfWriter()

        # Loop semua halaman, pisahkan ke 2 writer
        for i in range(total_pages):
            if i in extracted_indices:
                writer_extracted.add_page(reader.pages[i])
            else:
                writer_remaining.add_page(reader.pages[i])

        # Simpan file ke direktori sementara
        path_extracted = os.path.join(temp_dir, "extracted_pages.pdf")
        path_remaining = os.path.join(temp_dir, "remaining_pages.pdf")
        
        # Hanya simpan file jika berisi halaman
        if len(writer_extracted.pages) > 0:
            with open(path_extracted, "wb") as f_ext:
                writer_extracted.write(f_ext)
                
        if len(writer_remaining.pages) > 0:
            with open(path_remaining, "wb") as f_rem:
                writer_remaining.write(f_rem)
        
        # Buat file ZIP di memori
        zip_io = BytesIO()
        with zipfile.ZipFile(zip_io, 'w') as zf:
            if os.path.exists(path_extracted):
                zf.write(path_extracted, arcname="halaman_ekstrak.pdf")
            if os.path.exists(path_remaining):
                zf.write(path_remaining, arcname="halaman_sisa.pdf")
        
        zip_io.seek(0)

        return StreamingResponse(
            zip_io,
            media_type="application/zip",
            headers={"Content-Disposition": f"attachment; filename=split_{file.filename}.zip"},
            background=BackgroundTask(cleanup_dir, path=temp_dir)
        )
    except HTTPException as e:
        if 'temp_dir' in locals() and os.path.exists(temp_dir):
            cleanup_dir(temp_dir)
        raise e # Tampilkan error spesifik dari parse_page_range
    except Exception as e:
        if 'temp_dir' in locals() and os.path.exists(temp_dir):
            cleanup_dir(temp_dir)
        raise HTTPException(500, f"Terjadi error saat memisah PDF: {e}")


@app.post("/rotate", summary="Rotasi halaman PDF")
async def rotate_pdf(
    file: UploadFile = File(..., description="File PDF yang akan dirotasi."),
    angle: int = Form(..., description="Sudut rotasi (hanya 90, 180, 270)")
):
    # (Kode tidak berubah)
    if file.content_type != "application/pdf":
        raise HTTPException(400, "Hanya file PDF yang diizinkan.")
    if angle not in [90, 180, 270]:
        raise HTTPException(400, "Sudut rotasi harus 90, 180, or 270.")
    try:
        pdf_bytes = await file.read()
        reader = PdfReader(BytesIO(pdf_bytes))
        if reader.is_encrypted:
            raise HTTPException(400, "File PDF terenkripsi. Buka sandi terlebih dahulu.")
        writer = PdfWriter()
        for page in reader.pages:
            page.rotate(angle)
            writer.add_page(page)
        output_io = BytesIO()
        writer.write(output_io)
        output_io.seek(0)
        return StreamingResponse(
            output_io,
            media_type="application/pdf",
            headers={"Content-Disposition": f"attachment; filename=rotated_{file.filename}"}
        )
    except Exception as e:
        raise HTTPException(500, f"Terjadi error: {e}")


# --- FITUR BARU ---
@app.post("/to-powerpoint", summary="Konversi PDF ke PowerPoint (.pptx)")
async def pdf_to_powerpoint(file: UploadFile = File(..., description="File PDF yang akan dikonversi")):
    """
    Mengkonversi PDF ke PowerPoint.
    CATATAN: Setiap halaman PDF akan menjadi GAMBAR di setiap slide.
    Teks tidak dapat diedit.
    """
    if file.content_type != "application/pdf":
        raise HTTPException(400, "Hanya file PDF yang diizinkan.")

    try:
        # Baca PDF dari bytes
        pdf_bytes = await file.read()
        
        # Konversi PDF ke list gambar (Pillow)
        # Ingat, ini membutuhkan POPOPPLER
        images = convert_from_bytes(pdf_bytes)
        
        prs = Presentation()
        # Dapatkan ukuran slide default (landscape 10x7.5 inch)
        slide_width = prs.slide_width
        slide_height = prs.slide_height

        for img in images:
            # Layout 6 adalah layout kosong (blank)
            blank_slide_layout = prs.slide_layouts[6] 
            slide = prs.slides.add_slide(blank_slide_layout)

            # Simpan gambar dari Pillow ke buffer memori
            img_io = BytesIO()
            img.save(img_io, format='PNG')
            img_io.seek(0)
            
            # Tambahkan gambar, paskan ke tinggi slide
            pic = slide.shapes.add_picture(img_io, Inches(0), Inches(0), height=slide_height)
            
            # Atur posisi gambar agar di tengah (horizontal)
            pic.left = int((slide_width - pic.width) / 2)
            pic.top = 0

        # Simpan presentasi ke buffer memori
        output_io = BytesIO()
        prs.save(output_io)
        output_io.seek(0)

        # Kirim file .pptx
        return StreamingResponse(
            output_io,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": f"attachment; filename={os.path.splitext(file.filename)[0]}.pptx"}
        )

    except Exception as e:
        raise HTTPException(500, f"Terjadi error saat konversi ke PPTX: {e}. Pastikan Poppler terinstal.")


@app.post("/to-excel", summary="Konversi tabel PDF ke Excel (termasuk gambar)")
async def pdf_to_excel(
    file: UploadFile = File(..., description="File PDF yang akan dikonversi."),
    flavor: str = Form("lattice", description="Metode ekstraksi: 'lattice' (untuk tabel bergaris) atau 'stream' (tanpa garis).")
):
    """
    Mengekstrak tabel dari PDF dan menyimpannya sebagai file Excel.
    Akan mencoba mengekstrak gambar yang ada di dalam sel.
    
    MEMBUTUHKAN GHOSTSCRIPT terinstal di server.
    """
    if file.content_type != "application/pdf":
        raise HTTPException(400, "Hanya file PDF yang diizinkan.")
    if flavor not in ['lattice', 'stream']:
        raise HTTPException(400, "Flavor harus 'lattice' atau 'stream'.")

    temp_pdf_path = None
    output_path = None
    pdf_doc = None
    
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as temp_pdf:
            content = await file.read()
            temp_pdf.write(content)
            temp_pdf_path = temp_pdf.name
        
        output_path = tempfile.mktemp(suffix=".xlsx") 
        
        print("\n--- [DEBUG] Memulai Ekstraksi Camelot ---")
        tables = camelot.read_pdf(temp_pdf_path, pages='all', flavor=flavor)
        print(f"--- [DEBUG] Camelot Selesai. Ditemukan {tables.n} tabel. ---")

        if tables.n == 0:
            raise HTTPException(404, "Tidak ada tabel yang ditemukan di PDF ini.")

        pdf_doc = fitz.open(temp_pdf_path) 
        print("--- [DEBUG] PyMuPDF (fitz) berhasil membuka PDF. ---")

        with pd.ExcelWriter(output_path, engine='openpyxl') as writer:
            for i, table in enumerate(tables):
                sheet_name = f'Tabel {i+1}'
                table.df.to_excel(
                    writer, 
                    sheet_name=sheet_name, 
                    header=False, 
                    index=False
                )
                
                ws = writer.sheets[sheet_name]
                page_num = table.page - 1 
                page = pdf_doc.load_page(page_num)
                
                page_height = page.rect.height
                
                print(f"\n--- [DEBUG] Memproses Tabel {i+1} di Halaman PDF {page_num+1} (Tinggi: {page_height}pt) ---")
                
                images_on_page = page.get_images(full=True)
                print(f"--- [DEBUG] Ditemukan {len(images_on_page)} gambar di halaman ini. ---")

                # [BARU] Set untuk melacak xref gambar yang sudah digunakan
                used_image_xrefs = set()

                for r_idx in range(len(table.rows)):
                    for c_idx in range(len(table.cols)):
                        cell_text = table.df.iloc[r_idx, c_idx]
                        if not cell_text: # Hanya proses sel yang kosong
                            cell_coords = table.cells[r_idx][c_idx]
                            
                            new_y1 = page_height - cell_coords.y2
                            new_y2 = page_height - cell_coords.y1
                            
                            cell_bbox = fitz.Rect(cell_coords.x1, new_y1, cell_coords.x2, new_y2)
                            
                            print(f"\n[DEBUG] Memeriksa Sel Kosong [{r_idx},{c_idx}]")
                            print(f"  > Koord. Asli (Camelot, Bawah): ({cell_coords.x1}, {cell_coords.y1}, {cell_coords.x2}, {cell_coords.y2})")
                            print(f"  > Koord. Baru (PyMuPDF, Atas): {cell_bbox}")
                            
                            # [BARU] Flag untuk menandai apakah gambar sudah ditemukan untuk sel ini
                            image_found_for_cell = False

                            for img_info in images_on_page:
                                xref = img_info[0]
                                
                                # [BARU] Lewati gambar jika sudah digunakan
                                if xref in used_image_xrefs:
                                    print(f"  > Gambar (xref:{xref}) sudah digunakan. Lewati.")
                                    continue
                                
                                rects = page.get_image_rects(xref)
                                for r in rects:
                                    img_bbox = fitz.Rect(r)
                                    print(f"  > Membandingkan dengan Gbr (xref:{xref}) di Koordinat: {img_bbox}")
                                    
                                    if cell_bbox.intersects(img_bbox):
                                        print(f"  [BERHASIL!] Gambar BERSINGGUNGAN dengan sel. Mencoba menyisipkan...")
                                        
                                        base_image = pdf_doc.extract_image(xref)
                                        image_bytes = base_image["image"]
                                        
                                        try:
                                            img_data = BytesIO(image_bytes)
                                            excel_img = OpenPyXLImage(img_data)
                                            cell_id = f"{get_column_letter(c_idx + 1)}{r_idx + 1}"
                                            
                                            ws.row_dimensions[r_idx + 1].height = 70
                                            ws.column_dimensions[get_column_letter(c_idx + 1)].width = 15
                                            excel_img.height = 80
                                            excel_img.width = 80
                                            
                                            ws.add_image(excel_img, cell_id)
                                            print(f"  [SUKSES] Gambar disisipkan ke sel {cell_id}")
                                            
                                            used_image_xrefs.add(xref) # [BARU] Tandai gambar ini sudah digunakan
                                            image_found_for_cell = True # [BARU] Set flag
                                            break # Keluar dari loop 'r in rects'
                                        except Exception as e:
                                            print(f"  [ERROR SISPKA] Gagal memuat/menyisipkan gambar: {e}")
                                    else:
                                        print(f"  [INFO] Gambar tidak bersinggungan.")
                                
                                if image_found_for_cell: # [BARU] Jika sudah ketemu gambar untuk sel ini
                                    break # Keluar dari loop 'img_info in images_on_page'
                            # Akhir loop gambar
                    # Akhir loop kolom
                # Akhir loop baris
            # Akhir loop tabel
        
        print("\n--- [DEBUG] Menutup dokumen PDF. ---")
        pdf_doc.close()
        pdf_doc = None 
        
        os.remove(temp_pdf_path) 
        print("--- [DEBUG] File PDF sementara dihapus. ---")

        return FileResponse(
            path=output_path,
            media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            filename=f"{os.path.splitext(file.filename)[0]}.xlsx",
            background=BackgroundTask(cleanup_file, path=output_path)
        )

    except Exception as e:
        print(f"\n--- [DEBUG] Terjadi ERROR Global ---")
        print(str(e))
        if pdf_doc:
            pdf_doc.close()
        if temp_pdf_path and os.path.exists(temp_pdf_path):
            os.remove(temp_pdf_path)
        if output_path and os.path.exists(output_path):
            os.remove(output_path)
        
        if isinstance(e, HTTPException):
            raise e
        
        raise HTTPException(500, f"Terjadi error saat konversi ke Excel (dengan gambar): {e}. Pastikan Ghostscript terinstal.")

@app.post("/delete-pages", summary="Hapus halaman PDF berdasarkan rentang")
async def delete_pages(
    file: UploadFile = File(..., description="File PDF yang akan diproses."),
    page_range: str = Form(..., description="Halaman yang akan dihapus (cth: '2-5' atau '1, 3')")
):
    """
    Menghapus halaman berdasarkan rentang (cth: '1, 3, 5-7')
    dan mengembalikan PDF dengan sisa halaman.
    """
    if file.content_type != "application/pdf":
        raise HTTPException(400, "Hanya file PDF yang diizinkan.")

    try:
        pdf_bytes = await file.read()
        reader = PdfReader(BytesIO(pdf_bytes))

        if reader.is_encrypted:
            raise HTTPException(400, "File PDF terenkripsi. Buka sandi terlebih dahulu.")
        
        total_pages = len(reader.pages)
        
        # Gunakan helper yang sama dengan '/split'
        # untuk mendapatkan set halaman (0-indexed) yang akan DIHAPUS
        indices_to_delete = parse_page_range(page_range, total_pages)
        
        writer = PdfWriter()

        # Loop semua halaman, tambahkan HANYA jika TIDAK ADA di set hapus
        for i in range(total_pages):
            if i not in indices_to_delete:
                writer.add_page(reader.pages[i])
        
        if len(writer.pages) == 0:
            raise HTTPException(400, "Tidak ada halaman tersisa setelah penghapusan.")

        # Simpan ke memori
        output_io = BytesIO()
        writer.write(output_io)
        output_io.seek(0)
        
        return StreamingResponse(
            output_io,
            media_type="application/pdf",
            headers={"Content-Disposition": f"attachment; filename=deleted_{file.filename}"}
        )
    except HTTPException as e:
        raise e
    except Exception as e:
        raise HTTPException(500, f"Terjadi error saat menghapus halaman: {e}")


@app.post("/arrange-pages", summary="Atur ulang urutan halaman PDF")
async def arrange_pages(
    file: UploadFile = File(..., description="File PDF yang akan diatur."),
    new_order: str = Form(..., description="Urutan halaman baru, dipisah koma (cth: '3,1,2,4')")
):
    """
    Mengatur ulang halaman PDF berdasarkan urutan 1-indexed yang diberikan.
    Jumlah halaman di urutan baru harus sama dengan total halaman PDF.
    """
    if file.content_type != "application/pdf":
        raise HTTPException(400, "Hanya file PDF yang diizinkan.")

    try:
        pdf_bytes = await file.read()
        reader = PdfReader(BytesIO(pdf_bytes))

        if reader.is_encrypted:
            raise HTTPException(400, "File PDF terenkripsi. Buka sandi terlebih dahulu.")
        
        total_pages = len(reader.pages)
        
        # Parsing input 'new_order'
        try:
            # Ubah string "3,1,2,4" -> list [2, 0, 1, 3] (0-indexed)
            order_indices = [int(p.strip()) - 1 for p in new_order.split(',')]
        except ValueError:
            raise HTTPException(400, "Format 'new_order' tidak valid. Gunakan angka dipisah koma.")

        # Validasi
        if len(order_indices) != total_pages:
            raise HTTPException(400, f"Jumlah halaman di 'new_order' ({len(order_indices)}) tidak cocok dengan total halaman PDF ({total_pages}).")
        if not all(0 <= i < total_pages for i in order_indices):
            raise HTTPException(400, "Urutan halaman tidak valid (angka di luar rentang).")
        if len(set(order_indices)) != total_pages:
            raise HTTPException(400, "Urutan halaman tidak boleh ada duplikat.")

        writer = PdfWriter()
        
        # Tambahkan halaman sesuai urutan baru
        for index in order_indices:
            writer.add_page(reader.pages[index])

        output_io = BytesIO()
        writer.write(output_io)
        output_io.seek(0)
        
        return StreamingResponse(
            output_io,
            media_type="application/pdf",
            headers={"Content-Disposition": f"attachment; filename=arranged_{file.filename}"}
        )
    except HTTPException as e:
        raise e
    except Exception as e:
        raise HTTPException(500, f"Terjadi error saat mengatur halaman: {e}")


@app.post("/add-signature", summary="Tambahkan gambar tanda tangan ke PDF")
async def add_signature(
    file: UploadFile = File(..., description="File PDF utama."),
    signature_image: UploadFile = File(..., description="File gambar .png tanda tangan."),
    page_number: int = Form(1, description="Nomor halaman (1-indexed) untuk tanda tangan."),
    x_pos: int = Form(50, description="Posisi X (dari kiri) dalam poin (pt)."),
    y_pos: int = Form(50, description="Posisi Y (dari BAWAH) dalam poin (pt)."),
    width: int = Form(150, description="Lebar gambar tanda tangan dalam poin (pt).")
):
    """
    Menambahkan gambar (seperti tanda tangan) ke halaman PDF
    pada koordinat yang ditentukan.
    CATATAN: (0, 0) adalah pojok KIRI BAWAH.
    """
    if file.content_type != "application/pdf":
        raise HTTPException(400, "Hanya file PDF yang diizinkan.")
    if signature_image.content_type not in ["image/png", "image/jpeg"]:
        raise HTTPException(400, "File tanda tangan harus .png atau .jpg.")

    try:
        pdf_bytes = await file.read()
        reader = PdfReader(BytesIO(pdf_bytes))

        if reader.is_encrypted:
            raise HTTPException(400, "File PDF terenkripsi. Buka sandi terlebih dahulu.")
        
        page_index = page_number - 1
        if not (0 <= page_index < len(reader.pages)):
            raise HTTPException(400, "Nomor halaman tidak valid.")
            
        # --- [LOGIKA BARU DIMULAI DI SINI] ---
        
        # 1. Baca gambar tanda tangan
        sig_bytes = await signature_image.read()
        sig_io = BytesIO(sig_bytes)
        sig_pil_img = Image.open(sig_io)
        
        # Dapatkan rasio aspek untuk menghitung tinggi
        img_width, img_height = sig_pil_img.size
        aspect_ratio = img_height / img_width
        height = int(width * aspect_ratio) # Hitung tinggi otomatis

        # 2. Buat "Stempel" PDF di memori
        stamp_io = BytesIO()
        
        # Ambil ukuran halaman target agar stempel pas
        target_page_box = reader.pages[page_index].mediabox
        page_width = target_page_box.width
        page_height = target_page_box.height

        # Buat kanvas reportlab
        c = canvas.Canvas(stamp_io, pagesize=(page_width, page_height))
        
        # Gambar tanda tangan ke kanvas di posisi X, Y
        # (Reportlab dan pypdf sama-sama pakai Kiri-Bawah sebagai 0,0)
        c.drawImage(
            ImageReader(sig_io), # Gunakan ImageReader untuk BytesIO
            x_pos, 
            y_pos, 
            width=width, 
            height=height, 
            mask='auto' # Penting untuk transparansi PNG
        )
        c.save() # Simpan PDF stempel
        
        # 3. Baca stempel PDF yang baru dibuat
        stamp_io.seek(0)
        stamp_reader = PdfReader(stamp_io)
        stamp_page = stamp_reader.pages[0]

        # 4. Gabungkan stempel dengan halaman PDF asli
        writer = PdfWriter()
        for i in range(len(reader.pages)):
            page = reader.pages[i]
            
            # Jika ini halaman target, gabungkan (overlay) dengan stempel
            if i == page_index:
                page.merge_page(stamp_page)
            
            writer.add_page(page)
        # --- [LOGIKA BARU SELESAI] ---

        output_io = BytesIO()
        writer.write(output_io)
        output_io.seek(0)
        
        return StreamingResponse(
            output_io,
            media_type="application/pdf",
            headers={"Content-Disposition": f"attachment; filename=signed_{file.filename}"}
        )
    except HTTPException as e:
        raise e
    except Exception as e:
        raise HTTPException(500, f"Terjadi error saat menambah tanda tangan: {e}")

        
# --- Jalankan Server ---
if __name__ == "__main__":
    uvicorn.run("convert_pdf:app", host="0.0.0.0", port=8000, reload=True)